#!/usr/bin/env python3
"""
Conference one-pager generator
==============================

Turns the transcript of a single conference talk + the slides that were
presented into a tidy, web-ready one-pager (HTML) and a matching PDF.

Pipeline (4 stages):
  1. INGEST      transcript (.srt / .txt) and slides (.pptx / .png folder)
  2. SYNTHESISE  one Gemini call -> structured JSON (a OnePager object)
  3. RENDER      Jinja2 template fills the JSON -> a single .html file
  4. EXPORT      that exact .html is printed to .pdf  (one source of truth)

Usage:
  export GEMINI_API_KEY=...
  python onepager.py --transcript talk.srt --slides deck.pptx \
                     --conference "DevConf 2026" --date "2026-05-20" --out talk

  # when you only have screenshots of the slides:
  python onepager.py --transcript talk.txt --slides ./slide_pngs/ --out talk
"""
from __future__ import annotations

import argparse
import base64
import io
import json
import os
import random
import re
import sys
import time
import tomllib
import urllib.error
import urllib.request
from datetime import datetime
from dataclasses import dataclass, field
from pathlib import Path

import srt
from pptx import Presentation
from PIL import Image
from jinja2 import Environment, FileSystemLoader, select_autoescape
from pydantic import BaseModel
from google import genai
from google.genai import errors, types

MODEL = "gemini-2.5-flash"
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}
IG_WIDTH, IG_HEIGHT = 1080, 1350          # Instagram portrait post, 4:5 ratio


# ======================================================================
# 1. DATA SHAPES
# ----------------------------------------------------------------------
# The OnePager model is doing double duty: it is both (a) the JSON schema
# we hand to Gemini so the model is *forced* to return well-formed,
# predictable data, and (b) the object the template renders. Defining the
# shape once removes a whole class of "the model returned something my
# template didn't expect" bugs.
# ======================================================================
class Section(BaseModel):
    heading: str
    paragraphs: list[str]


class FeaturedSlide(BaseModel):
    """A slide the model judged worth embedding in the article."""
    slide_number: int              # 1-based index into the slide images given
    caption: str                   # one short line on what the slide shows


class Instagram(BaseModel):
    """Social copy that promotes the one-pager on Instagram.

    Note: the post's bullet points are NOT stored here -- they are taken
    directly from the article's `key_takeaways`, so the post and the page
    always say the same thing.
    """
    headline: str                  # punchy line for the image card (<= 8 words)
    caption: str                   # the IG caption prose (no hashtags, no URL)
    hashtags: list[str]            # 4-8 tags, each starting with '#'


class OnePager(BaseModel):
    title: str
    subtitle: str
    speakers: list[str]
    dek: str                       # the lead-in / hook paragraph
    sections: list[Section]
    key_takeaways: list[str]
    pull_quote: str                # one near-verbatim line from the talk
    topics: list[str]              # tags
    featured_slides: list[FeaturedSlide]   # 0-2 slides to embed as figures
    instagram: Instagram           # social copy promoting the one-pager


@dataclass
class Slides:
    """Slide content for the model. EITHER field may be populated, or BOTH.

    text   -- extracted from .pptx (titles, bullets, speaker notes)
    images -- PIL images from .png screenshots
    A single talk can supply a .pptx AND screenshots; both get merged here.
    """
    text: str = ""
    images: list = field(default_factory=list)   # list[PIL.Image]


# ======================================================================
# 2. INGEST -- TRANSCRIPT
# ----------------------------------------------------------------------
# .srt has precise, structured timestamps; .txt has looser, free-form
# ones. For a one-pager we do NOT need timestamps at all (we are not
# aligning slides to spoken segments in v1), so both paths converge on
# the same thing: clean plain text. We just have to strip the timestamps
# tolerantly because the .txt format is not guaranteed.
# ======================================================================
_BRACKET_TS = re.compile(r'^\s*[\[(]\s*\d{1,2}:\d{2}(?::\d{2})?(?:[.,]\d+)?\s*[\])]\s*')
_BARE_TS = re.compile(r'^\s*\d{1,2}:\d{2}(?::\d{2})?(?:[.,]\d+)?\s+')
_TAG = re.compile(r'<[^>]+>')                       # <i>, <b>, <v Speaker> ...


def _strip_timestamps(line: str) -> str:
    """Remove a leading timestamp token, whatever common shape it takes."""
    prev = None
    while prev != line:                  # loop: handles "[00:01] 00:01 text"
        prev = line
        line = _BRACKET_TS.sub('', line)
        line = _BARE_TS.sub('', line)
    return line.strip()


def read_srt(path: Path) -> str:
    """Parse .srt into one clean paragraph, dropping repeated caption rows."""
    raw = path.read_text(encoding="utf-8", errors="ignore")
    out: list[str] = []
    last = None
    for cue in srt.parse(raw):
        text = _TAG.sub('', cue.content).replace('\n', ' ').strip()
        if text and text != last:        # rolling captions repeat lines
            out.append(text)
            last = text
    return ' '.join(out)


def read_txt_transcript(path: Path) -> str:
    """Parse a plain-text transcript, stripping leading timestamps."""
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    cleaned = (_strip_timestamps(_TAG.sub('', ln)) for ln in lines)
    return '\n'.join(ln for ln in cleaned if ln)


def load_transcript(path: Path) -> str:
    if not path.exists():
        sys.exit(f"Transcript not found: {path}")
    if path.suffix.lower() == ".srt":
        return read_srt(path)
    return read_txt_transcript(path)


# ======================================================================
# 3. INGEST -- SLIDES
# ----------------------------------------------------------------------
# .pptx  -> pull out title / body / speaker-notes as TEXT (free, lossless
#           for text; speaker notes are often the verbatim script -> gold)
# .png   -> keep the IMAGES and let the multimodal model read them. We do
#           NOT OCR: OCR discards charts, layout and emphasis -- the very
#           things that make a slide a slide.
# --slides accepts ONE OR MORE inputs, so the same talk can supply a
# .pptx and a folder of screenshots; load_slides() merges them.
# ======================================================================
def read_pptx(path: Path) -> Slides:
    prs = Presentation(str(path))
    blocks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if not txt:
                continue
            if shape == slide.shapes.title:
                title = txt
            else:
                body.append(txt)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        block = [f"--- Slide {i}: {title or '(untitled)'} ---"]
        if body:
            block.append("\n".join(body))
        if notes:
            block.append(f"[Speaker notes] {notes}")
        blocks.append("\n".join(block))
    return Slides(text="\n\n".join(blocks))


def _natural_key(p: Path):
    """Sort slide_2.png before slide_10.png (plain sort gets this wrong)."""
    return [int(t) if t.isdigit() else t.lower()
            for t in re.split(r'(\d+)', p.name)]


def read_slide_images(path: Path) -> Slides:
    if path.is_dir():
        files = sorted((f for f in path.iterdir()
                        if f.suffix.lower() in IMAGE_EXTS), key=_natural_key)
    else:
        files = [path]
    if not files:
        sys.exit(f"No slide images found in: {path}")
    images = [Image.open(f).convert("RGB") for f in files]
    return Slides(images=images)


def _load_one_slides_input(path: Path) -> Slides:
    """Ingest a single slides input (one .pptx, one image, or one folder)."""
    if not path.exists():
        sys.exit(f"Slides input not found: {path}")
    if path.is_dir():
        return read_slide_images(path)
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return read_pptx(path)
    if suffix in IMAGE_EXTS:
        return read_slide_images(path)
    sys.exit(f"Unsupported slides input: {path} (expected .pptx, image, or folder)")


def load_slides(paths: list[Path]) -> Slides:
    """Ingest one OR MORE slide inputs and merge them into a single Slides.

    Merge rule: .pptx text blocks are concatenated; image lists are
    extended in order. This is what lets a single talk feed both a deck
    and a folder of screenshots in one run.
    """
    merged = Slides()
    for path in paths:
        part = _load_one_slides_input(path)
        if part.text:
            merged.text = "\n\n".join(t for t in (merged.text, part.text) if t)
        merged.images.extend(part.images)
    if not merged.text and not merged.images:
        sys.exit("No usable slide content found in the given input(s).")
    return merged


# ======================================================================
# 4. SYNTHESISE  -- the single Gemini call
# ----------------------------------------------------------------------
# Two things make this reliable:
#  - response_schema=OnePager  -> the model MUST return our exact shape.
#  - the system prompt hammers GROUNDING: only use what was presented.
#    "based solely on what was presented" is a hallucination constraint,
#    so it has to live in the prompt explicitly, not be hoped for.
# ======================================================================
SYSTEM_PROMPT = """\
You are an editor producing a concise, engaging one-page web article that
summarises a SINGLE conference talk for a general professional audience.

HARD RULES — follow exactly:
- Use ONLY information contained in the supplied transcript and slides. Do
  not add outside facts, statistics, names, history, or context. If it was
  not presented, it does not go in.
- NO EXAGGERATION. This is a medical/health context, so accuracy is critical.
  Do not overstate findings, hype results, or imply more certainty than the
  speakers expressed. Never claim or imply a cure, breakthrough, or guaranteed
  outcome unless the speakers explicitly and unambiguously said so. Hedge as
  the speakers hedged: if they said a treatment "may help" or is "in trials",
  do not upgrade that to "works" or "is available". Prefer cautious, precise
  wording over dramatic wording everywhere -- in the article AND the Instagram
  copy.
- `pull_quote` MUST be a near-verbatim sentence actually spoken in the
  transcript. Never invent or paraphrase a quote into existence.
- Keep it genuinely ONE PAGE: 400-600 words total across `dek` + `sections`.
- Write clear, lively, specific prose. No marketing fluff, no throat-clearing
  ("in today's fast-paced world"), no empty superlatives.
- Produce 3-4 `sections`, each a real sub-topic of the talk; 3-5
  `key_takeaways` ORDERED most important first (the first 3 are also used as
  the Instagram post's points, so they must stand on their own); 3-6 `topics`.
- `dek` is a single punchy lead paragraph (~40-60 words) that hooks the reader.
- `featured_slides`: pick the slide(s) most worth showing in the article --
  ones with charts, data, diagrams, or key figures. Prefer just ONE; add a
  second only if it clearly adds value. Skip title, agenda, and plain-text
  slides. If no slide images were provided, or none is visually informative,
  return an empty list. Each entry has the slide's 1-based number (as given
  in the SLIDE IMAGES section) and a caption of 12 words or fewer.
- `instagram`: short social copy promoting this one-pager.
  * `headline`: a punchy but accurate hook of 8 words or fewer for the post
    image. No hype, no overclaiming (see the NO EXAGGERATION rule).
  * `caption`: 2-4 natural sentences for the Instagram caption that hook the
    reader and make clear a full recap is now available. Do NOT include
    hashtags or any URL in this field.
  * `hashtags`: 4-8 relevant hashtags, each starting with '#'.
"""


def build_prompt(transcript: str, slides: Slides, meta: dict) -> str:
    parts = [
        f"Conference: {meta.get('conference') or '(unknown)'}",
        f"Date: {meta.get('date') or '(unknown)'}",
        "",
        "=== TRANSCRIPT ===",
        transcript,
        "",
    ]
    if slides.text:
        parts += ["=== SLIDE TEXT (titles, bullets, speaker notes) ===",
                  slides.text, ""]
    if slides.images:
        parts += ["=== SLIDE IMAGES ===",
                  f"{len(slides.images)} slide image(s) are attached below, in "
                  f"order. They are numbered 1 to {len(slides.images)}; refer to "
                  f"them by these numbers in `featured_slides`."]
    # When BOTH are present they describe the same deck -- say so explicitly
    # so the model treats them as one source, not two.
    if slides.text and slides.images:
        parts += ["",
                  "NOTE: the slide text and the slide images are the SAME "
                  "presentation. Use the text for exact wording and speaker "
                  "notes; use the images for visual content (charts, demos, "
                  "screenshots) that the text cannot capture. Do not treat "
                  "them as separate material or double-count points."]
    return "\n".join(parts)


def _generate_with_retry(client: genai.Client, *, model, contents, config,
                         max_attempts: int = 5, base_delay: float = 2.0):
    """Call Gemini, retrying ONLY transient failures with exponential backoff.

    RETRY:
      - errors.ServerError  -- any 5xx (e.g. 503 overloaded): server-side,
        temporary by nature.
      - errors.ClientError with code 429 -- rate-limited / quota: clears
        with time.
    DO NOT RETRY:
      - any other ClientError (400 malformed, 401/403 bad key, 404): the
        request itself is wrong, so an identical retry fails identically.

    Backoff: waits 2s, 4s, 8s, 16s ... doubling each attempt. The doubling
    gives an overloaded server real room to recover instead of being hit
    again immediately. The random `jitter` added on top stops many clients
    (or many talks in a batch) from retrying in lockstep and re-colliding.
    """
    last_error = None
    for attempt in range(1, max_attempts + 1):
        try:
            return client.models.generate_content(
                model=model, contents=contents, config=config)
        except errors.ServerError as e:                 # 5xx -> retry
            last_error = e
        except errors.ClientError as e:                 # 4xx -> retry only 429
            if getattr(e, "code", None) != 429:
                raise                                   # permanent: give up now
            last_error = e

        if attempt == max_attempts:
            raise last_error                            # out of attempts

        delay = base_delay * 2 ** (attempt - 1)         # 2, 4, 8, 16 ...
        delay += random.uniform(0, 0.3 * delay)         # jitter
        code = getattr(last_error, "code", "?")
        print(f"  API returned {code} (transient) -- "
              f"retry {attempt}/{max_attempts - 1} in {delay:.1f}s ...")
        time.sleep(delay)


def synthesize(client: genai.Client, transcript: str,
               slides: Slides, meta: dict) -> OnePager:
    # Gemini `contents` is a list of parts; strings and PIL images can be
    # mixed freely. Slide text (if any) is already inside the prompt string;
    # slide images (if any) are appended here. Either or both may be present.
    contents: list = [build_prompt(transcript, slides, meta)]
    if slides.images:
        contents.extend(slides.images)

    config = types.GenerateContentConfig(
        system_instruction=SYSTEM_PROMPT,
        response_mime_type="application/json",
        response_schema=OnePager,
        temperature=0.2,                # low: literal, sticks to transcript
    )
    response = _generate_with_retry(
        client, model=MODEL, contents=contents, config=config)

    if response.parsed is None:
        raise RuntimeError(
            "Model did not return valid structured output:\n" + (response.text or ""))
    return response.parsed


# ======================================================================
# 5. RENDER  -- HTML now, PDF from that same HTML
# ----------------------------------------------------------------------
# One source of truth: the template is the website one-pager AND the input
# to the PDF. Images (logo, featured slides) are embedded as base64 data
# URIs so the .html file is fully self-contained -- no sidecar image files
# to ship to the website or lose before the PDF step.
# ======================================================================
_MIME = {".svg": "image/svg+xml", ".png": "image/png", ".jpg": "image/jpeg",
         ".jpeg": "image/jpeg", ".webp": "image/webp"}


def file_to_data_uri(path: Path) -> str:
    """Embed any image file (incl. .svg) as a base64 data URI."""
    mime = _MIME.get(path.suffix.lower(), "application/octet-stream")
    b64 = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{b64}"


def image_to_data_uri(img: Image.Image) -> str:
    """Embed an in-memory PIL image (a slide screenshot) as a PNG data URI."""
    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    b64 = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/png;base64,{b64}"


def build_figures(data: OnePager, slides: Slides) -> list[dict]:
    """Resolve the model's chosen slide numbers into embeddable figures.

    The model only returns slide *numbers*; the actual pixels live in
    `slides.images`. This maps one to the other, capping at 2, and skipping
    anything the model picked that we cannot honour (no images supplied, or
    an out-of-range number).
    """
    figures: list[dict] = []
    if not slides.images:
        if data.featured_slides:
            print("  note: model selected featured slides, but no slide images "
                  "were supplied -- skipping embedded figures.")
        return figures
    count = len(slides.images)
    for fs in data.featured_slides[:2]:                    # hard cap at 2
        n = fs.slide_number
        if 1 <= n <= count:
            figures.append({"uri": image_to_data_uri(slides.images[n - 1]),
                            "caption": fs.caption})
        else:
            print(f"  note: featured slide #{n} is out of range "
                  f"(1-{count}) -- skipped.")
    return figures


# ----------------------------------------------------------------------
# YouTube thumbnail
# ----------------------------------------------------------------------
# A YouTube link comes in several shapes -- youtube.com/watch?v=ID,
# youtu.be/ID, /embed/ID, /shorts/ID. We pull the 11-char video ID out of
# whichever shape we got, then fetch YouTube's own thumbnail for it.
_YT_ID = re.compile(
    r'(?:v=|/embed/|/shorts/|youtu\.be/)([A-Za-z0-9_-]{11})')


def youtube_video_id(url: str) -> str | None:
    """Extract the 11-character video ID from any common YouTube URL form."""
    match = _YT_ID.search(url)
    return match.group(1) if match else None


def fetch_youtube_thumbnail(url: str) -> str:
    """Return a data URI for the video's thumbnail, or '' on any failure.

    YouTube serves thumbnails at predictable URLs. The catch: `maxresdefault`
    (1280x720) only exists if the uploader's source was that big -- for many
    videos it 404s. `hqdefault` (480x360) is generated for EVERY video. So we
    try the high-res one first and fall back to the one that always exists.
    Any network problem just returns '' -- a missing thumbnail must never
    break the page.
    """
    vid = youtube_video_id(url)
    if not vid:
        print(f"  note: could not find a video ID in '{url}' -- no thumbnail.")
        return ""
    candidates = [
        f"https://img.youtube.com/vi/{vid}/maxresdefault.jpg",   # may 404
        f"https://img.youtube.com/vi/{vid}/hqdefault.jpg",       # always exists
    ]
    for thumb_url in candidates:
        try:
            req = urllib.request.Request(
                thumb_url, headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = resp.read()
            if len(data) < 1500:        # YouTube's "no image" placeholder is tiny
                continue
            b64 = base64.b64encode(data).decode("ascii")
            return f"data:image/jpeg;base64,{b64}"
        except (urllib.error.URLError, TimeoutError, OSError):
            continue
    print("  note: could not fetch a YouTube thumbnail -- using a text link.")
    return ""


def resolve_thumbnail(youtube_url: str, custom: Path | None) -> str:
    """Pick the video thumbnail: a custom image if given, else auto-fetch."""
    if custom:
        return file_to_data_uri(custom)        # explicit override wins
    if youtube_url:
        return fetch_youtube_thumbnail(youtube_url)
    return ""


# ----------------------------------------------------------------------
# SEO: JSON-LD structured data
# ----------------------------------------------------------------------
# Google reads JSON-LD anywhere on the page, so we emit it into the BODY
# of the template -- which means it survives being pasted into WordPress
# even when no SEO plugin is installed. We mark up an Article always, and
# nest a VideoObject ONLY when we have a valid date (Google REQUIRES an
# uploadDate for VideoObject; a missing/guessed date invalidates it). We
# never fabricate fields -- absent data means a smaller, still-valid block.

# date formats a human might plausibly type into --date or the config
_DATE_FORMATS = [
    "%Y-%m-%d",        # 2024-05-30
    "%d/%m/%Y",        # 30/05/2024
    "%m/%d/%Y",        # 05/30/2024
    "%B %d, %Y",       # May 30, 2024
    "%b %d, %Y",       # May 30, 2024 (abbrev)
    "%d %B %Y",        # 30 May 2024
    "%B %Y",           # May 2024  (day unknown)
    "%Y",              # 2024      (year only)
]


def parse_iso_date(raw: str) -> str | None:
    """Return an ISO 8601 date string (YYYY-MM-DD) or None if unparseable.

    Google rejects non-ISO dates outright, so we convert known human
    formats and refuse anything we cannot parse with confidence -- a wrong
    date is worse than no date.
    """
    raw = (raw or "").strip()
    if not raw:
        return None
    for fmt in _DATE_FORMATS:
        try:
            dt = datetime.strptime(raw, fmt)
        except ValueError:
            continue
        # partial formats: pad sensibly so the result is still valid ISO
        if fmt == "%Y":
            return f"{dt.year:04d}-01-01"
        if fmt == "%B %Y":
            return f"{dt.year:04d}-{dt.month:02d}-01"
        return dt.strftime("%Y-%m-%d")
    return None


def build_jsonld(data: OnePager, meta: dict, page_url: str = "") -> str:
    """Build a JSON-LD structured-data block: Article (+ VideoObject if able).

    Returns a complete <script type="application/ld+json"> string, or '' if
    there is nothing worth emitting. Only fields backed by real page content
    are included -- never invented.
    """
    # description: reuse the article's lead paragraph, trimmed
    description = (data.dek or "").strip()
    if len(description) > 250:
        description = description[:247].rstrip() + "..."

    article: dict = {
        "@context": "https://schema.org",
        "@type": "Article",
        "headline": data.title[:110],        # schema caps headline at 110
        "description": description,
    }
    if data.speakers:
        article["author"] = [{"@type": "Person", "name": s}
                             for s in data.speakers]
    iso_date = parse_iso_date(meta.get("date", ""))
    if iso_date:
        article["datePublished"] = iso_date
    if page_url:
        article["mainEntityOfPage"] = {"@type": "WebPage", "@id": page_url}

    blocks = [article]

    # VideoObject -- only if Google's required trio can be satisfied:
    # name + uploadDate + (a video URL). No date -> no VideoObject.
    youtube = meta.get("youtube", "")
    if youtube and iso_date:
        video = {
            "@context": "https://schema.org",
            "@type": "VideoObject",
            "name": data.title[:110],
            "description": description,
            "uploadDate": iso_date,                 # Google REQUIRES this
            "embedUrl": youtube,
        }
        blocks.append(video)

    payload = blocks[0] if len(blocks) == 1 else blocks
    return ('<script type="application/ld+json">\n'
            + json.dumps(payload, indent=2, ensure_ascii=False)
            + '\n</script>')


def render_html(data: OnePager, meta: dict, template_dir: Path,
                logo_uri: str = "", figures: list[dict] | None = None,
                thumb_uri: str = "", page_url: str = "") -> str:
    env = Environment(
        loader=FileSystemLoader(str(template_dir)),
        autoescape=select_autoescape(["html"]),
    )
    template = env.get_template("template.html")
    jsonld = build_jsonld(data, meta, page_url)
    return template.render(d=data, meta=meta, logo_uri=logo_uri,
                           figures=figures or [], thumb_uri=thumb_uri,
                           jsonld=jsonld)


def html_to_pdf(html_path: Path, pdf_path: Path) -> None:
    """Print the HTML with a real browser engine for pixel-perfect parity.

    Requires a one-time:  playwright install chromium
    """
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page()
        page.goto(html_path.resolve().as_uri())
        page.wait_for_load_state("networkidle")   # let web fonts load
        page.pdf(
            path=str(pdf_path),
            format="A4",
            print_background=True,
            # margins are handled by the @page rule in the CSS:
            margin={"top": "0", "bottom": "0", "left": "0", "right": "0"},
        )
        browser.close()


def html_to_image(html: str, image_path: Path,
                  width: int, height: int) -> None:
    """Render an HTML string to a fixed-size PNG via a headless browser.

    Mirrors html_to_pdf, but `screenshot` instead of `pdf`: the IG template
    is a fixed pixel canvas, not a paged document. device_scale_factor=2
    renders at 2x so text stays crisp after Instagram's own compression.
    """
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": width, "height": height},
            device_scale_factor=2)
        page.set_content(html, wait_until="networkidle")   # waits for fonts
        page.wait_for_timeout(150)                         # final paint
        page.screenshot(path=str(image_path))
        browser.close()


# ======================================================================
# 6. INSTAGRAM POST  -- a second deliverable derived from the same data
# ----------------------------------------------------------------------
# The model already wrote the social copy into data.instagram (one call,
# one schema). Here we just render it two ways: the image card and the
# paste-ready caption text. Both read from the SAME structured fields.
# ======================================================================
def render_instagram_html(data: OnePager, meta: dict, template_dir: Path,
                          logo_uri: str = "", url: str = "") -> str:
    env = Environment(
        loader=FileSystemLoader(str(template_dir)),
        autoescape=select_autoescape(["html"]),
    )
    template = env.get_template("instagram_template.html")
    return template.render(d=data, meta=meta, logo_uri=logo_uri, url=url)


def build_caption(data: OnePager, url: str = "") -> str:
    """Assemble the paste-ready Instagram caption from the structured copy.

    The bullet points come from the article's top 3 key_takeaways, so the
    post and the one-pager always say the same thing.
    """
    ig = data.instagram
    lines = [ig.caption.strip(), ""]
    for point in data.key_takeaways[:3]:          # top 3 takeaways, scannable
        lines.append(f"\u2022 {point}")
    lines.append("")
    if url:
        lines.append(f"Read the full one-pager: {url}")
    else:
        lines.append("Read the full one-pager \u2014 link in bio.")
    lines.append("")
    lines.append(" ".join(ig.hashtags))
    return "\n".join(lines)


def _slugify(text: str) -> str:
    """Turn a title into a clean URL slug: 'Avastin & NF2!' -> 'avastin-nf2'."""
    text = text.lower()
    text = re.sub(r"[^a-z0-9]+", "-", text)     # non-alphanumerics -> hyphen
    return text.strip("-")[:60]                 # trim and cap length


def build_seo_sheet(data: OnePager, meta: dict, page_url: str = "") -> str:
    """A paste-ready SEO sheet for a WordPress plugin (Rank Math / Yoast).

    These are SUGGESTED values -- the human pastes them into the plugin's
    fields. Title and description are length-capped to Google's display
    limits so they don't get truncated in search results.
    """
    title = data.title.strip()
    # SEO title: Google shows ~60 chars; append the org name if it fits
    org = meta.get("conference", "").split(" - ")[0].strip()
    seo_title = title
    if org and org.lower() not in title.lower():
        candidate = f"{title} | {org}"
        seo_title = candidate if len(candidate) <= 60 else title
    seo_title = seo_title[:60]

    # meta description: ~155 chars, from the lead paragraph
    desc = " ".join((data.dek or "").split())
    if len(desc) > 155:
        desc = desc[:152].rstrip() + "..."

    slug = _slugify(title)
    keyword = (data.topics[0] if data.topics else title).strip()

    lines = [
        "SEO SHEET -- paste these into your WordPress SEO plugin",
        "(Rank Math or Yoast: edit the page -> find the SEO fields)",
        "=" * 58,
        "",
        f"SEO TITLE (<=60 chars, currently {len(seo_title)}):",
        f"  {seo_title}",
        "",
        f"META DESCRIPTION (<=155 chars, currently {len(desc)}):",
        f"  {desc}",
        "",
        "URL SLUG (the page's web address ending):",
        f"  {slug}",
        "",
        "FOCUS KEYWORD (the main term this page should rank for):",
        f"  {keyword}",
    ]
    if page_url:
        lines += ["", "CANONICAL URL:", f"  {page_url}"]
    lines += [
        "",
        "-" * 58,
        "Note: these are suggestions. Review them before publishing --",
        "especially the focus keyword, which should match what people",
        "actually search for.",
    ]
    return "\n".join(lines)


# ======================================================================
# 7. ORCHESTRATION
# ======================================================================
DEFAULT_CONFIG = "onepager.toml"

# Config keys that name a file/path -- these get wrapped in Path() so the
# rest of the program sees the same type whether a value came from a flag
# (argparse already does type=Path) or from the TOML file (plain strings).
_PATH_KEYS = {"transcript", "logo", "thumbnail", "out"}
# 'slides' is special: always a LIST of Paths, never a single one.


def load_config(explicit: Path | None) -> dict:
    """Load a TOML config: the one named with --config, else onepager.toml.

    Returns {} when there is nothing to load. A bad TOML file is a hard
    error (the user clearly meant to use it); a simply-absent default file
    is silent (it is optional by design).
    """
    if explicit is not None:
        path = explicit
        if not path.exists():
            sys.exit(f"Config file not found: {path}")
    else:
        path = Path(DEFAULT_CONFIG)
        if not path.exists():
            return {}                       # no default file -- that's fine

    try:
        with path.open("rb") as fh:         # tomllib requires binary mode
            config = tomllib.load(fh)
    except tomllib.TOMLDecodeError as exc:
        sys.exit(f"Config file {path} is not valid TOML: {exc}")
    print(f"- loaded config from {path}")
    return config


def merge_config(args: argparse.Namespace, config: dict) -> None:
    """Layer the config UNDER the command line, in place on `args`.

    Precedence (low -> high): code default < config file < CLI flag.
    Rule: the config only fills a value the user did NOT pass on the CLI.
    We detect "did not pass" by the value still equalling argparse's
    default -- so an explicit flag always wins over the file.
    """
    defaults = {
        "transcript": None, "slides": None, "out": "one-pager",
        "logo": None, "conference": "", "date": "", "url": "",
        "youtube": "", "thumbnail": None, "pdf": False, "no_instagram": False,
        "no_seo": False,
    }
    for key, default in defaults.items():
        if getattr(args, key) != default:
            continue                        # user passed it on the CLI -- keep
        if key not in config:
            continue                        # not in the file either -- skip
        value = config[key]
        if key == "slides":
            items = value if isinstance(value, list) else [value]
            setattr(args, key, [Path(p) for p in items])
        elif key in _PATH_KEYS:
            setattr(args, key, Path(value))
        else:
            setattr(args, key, value)


def main() -> None:
    ap = argparse.ArgumentParser(
        description="Generate a one-pager from a conference talk.")
    ap.add_argument("--config", type=Path, default=None,
                    help=f"TOML config file (auto-loads {DEFAULT_CONFIG} if present)")
    ap.add_argument("--transcript", type=Path, default=None,
                    help=".srt or .txt transcript file")
    ap.add_argument("--slides", nargs="+", type=Path, default=None,
                    help="one or more slide inputs: a .pptx, a slide image, "
                         "a folder of PNGs -- or several of these together")
    ap.add_argument("--out", default="one-pager",
                    help="output path prefix (writes <prefix>.html and .pdf)")
    ap.add_argument("--logo", type=Path, default=None,
                    help="optional logo image (.svg/.png/.jpg) shown top-left")
    ap.add_argument("--conference", default="", help="conference name (metadata)")
    ap.add_argument("--date", default="", help="talk date (metadata)")
    ap.add_argument("--url", default="",
                    help="public URL where this recap is published; used as the "
                         "Instagram CTA link AND the page's canonical URL for SEO")
    ap.add_argument("--youtube", default="",
                    help="URL of the talk recording, linked at the foot of the page")
    ap.add_argument("--thumbnail", type=Path, default=None,
                    help="custom video thumbnail image; overrides the auto-fetched "
                         "YouTube thumbnail")
    ap.add_argument("--pdf", action="store_true",
                    help="also export a PDF (HTML is always written)")
    ap.add_argument("--no-instagram", action="store_true",
                    help="skip the Instagram post image and caption")
    ap.add_argument("--no-seo", action="store_true",
                    help="skip the SEO sheet (structured data stays in the HTML)")
    args = ap.parse_args()

    # layer the config file UNDER the CLI flags, then enforce the values
    # that are mandatory no matter where they came from
    merge_config(args, load_config(args.config))
    missing = [name for name in ("transcript", "slides")
               if getattr(args, name) is None]
    if missing:
        sys.exit("Missing required input(s): "
                 + ", ".join("--" + m for m in missing)
                 + " -- pass them as flags or set them in the config file.")

    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        sys.exit("Set GEMINI_API_KEY in your environment.")

    print("- reading transcript ...")
    transcript = load_transcript(args.transcript)
    if len(transcript.split()) < 50:
        print("  WARNING: transcript looks very short -- check the file.")

    print("- reading slides ...")
    slides = load_slides(args.slides)
    bits = []
    if slides.text:
        bits.append("text")
    if slides.images:
        bits.append(f"{len(slides.images)} image(s)")
    print(f"  slides: {' + '.join(bits)}")

    print(f"- calling {MODEL} ...")
    client = genai.Client(api_key=api_key)
    meta = {"conference": args.conference, "date": args.date,
            "youtube": args.youtube}
    data = synthesize(client, transcript, slides, meta)

    words = (len(data.dek.split())
             + sum(len(p.split()) for s in data.sections for p in s.paragraphs))
    print(f"  generated ~{words} words across {len(data.sections)} sections.")
    if words > 700:
        print("  WARNING: long for one page -- with figures this may overflow.")

    # resolve logo + the slides the model chose to feature
    logo_uri = file_to_data_uri(args.logo) if args.logo else ""
    figures = build_figures(data, slides)
    if figures:
        print(f"  embedding {len(figures)} featured slide(s).")

    # resolve the video thumbnail (custom image, else auto-fetch from YouTube)
    thumb_uri = resolve_thumbnail(args.youtube, args.thumbnail)
    if thumb_uri:
        print("  video thumbnail resolved.")

    template_dir = Path(__file__).parent
    html = render_html(data, meta, template_dir,
                       logo_uri=logo_uri, figures=figures,
                       thumb_uri=thumb_uri, page_url=args.url)
    html_path = Path(f"{args.out}.html")
    html_path.write_text(html, encoding="utf-8")
    print(f"- wrote {html_path}")

    if args.pdf:
        pdf_path = Path(f"{args.out}.pdf")
        html_to_pdf(html_path, pdf_path)
        print(f"- wrote {pdf_path}")

    if not args.no_seo:
        seo_path = Path(f"{args.out}_seo.txt")
        seo_path.write_text(build_seo_sheet(data, meta, args.url),
                            encoding="utf-8")
        print(f"- wrote {seo_path}")

    if not args.no_instagram:
        ig_html = render_instagram_html(data, meta, template_dir,
                                        logo_uri=logo_uri, url=args.url)
        ig_image = Path(f"{args.out}_instagram.png")
        html_to_image(ig_html, ig_image, IG_WIDTH, IG_HEIGHT)
        caption_path = Path(f"{args.out}_caption.txt")
        caption_path.write_text(build_caption(data, args.url), encoding="utf-8")
        print(f"- wrote {ig_image} and {caption_path}")

    print("done.")


if __name__ == "__main__":
    main()
